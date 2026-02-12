import streamlit as st
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re
import plotly.express as px
import plotly.graph_objects as go
from sklearn.metrics import jaccard_score
from scipy.spatial.distance import pdist, squareform

# ==========================================
# 1. EXACT TEXT EXTRACTION & PARSING ENGINE
# ==========================================

class PPTParser:
    def __init__(self, uploaded_file):
        self.uploaded_file = uploaded_file
        self.logic_rows = []
        self.risk_rows = []

    def _extract_text_exact_logic(self):
        """
        Uses the EXACT user-provided logic to process shapes/tables into a clean text stream.
        """
        prs = Presentation(self.uploaded_file)
        slide_contents = {}

        for i, slide in enumerate(prs.slides, start=1):
            slide_lines = []
            
            # 1. Collect shapes
            shapes_on_slide = []
            for shape in slide.shapes:
                if shape.has_text_frame or shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    shapes_on_slide.append(shape)

            # 2. Sort Visually (Top-to-Bottom, Left-to-Right)
            shapes_on_slide.sort(key=lambda s: (int(s.top), s.left))

            # 3. Extract Text
            for shape in shapes_on_slide:
                
                # --- Text Boxes ---
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        # Fix Soft Returns
                        raw_text = paragraph.text.replace('\x0b', '\n')
                        sub_lines = raw_text.split('\n')
                        
                        for line in sub_lines:
                            clean_line = line.strip()
                            # Remove separator lines (==== or ----)
                            if clean_line and not re.match(r'^[-=]{3,}$', clean_line):
                                slide_lines.append(clean_line)

                # --- Tables ---
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    for row in shape.table.rows:
                        row_cells = []
                        for cell in row.cells:
                            if cell.text_frame:
                                cell_text = " ".join([p.text.strip() for p in cell.text_frame.paragraphs if p.text.strip()])
                                if cell_text:
                                    row_cells.append(cell_text)
                        if row_cells:
                            slide_lines.append(" | ".join(row_cells))

            # Join slide content with Double Newline for clear separation
            slide_contents[i] = "\n\n".join(slide_lines)
            
        return slide_contents

    def parse(self):
        """
        Main Parsing Loop. Uses the extracted text to build the database.
        Maintains context (Metadata) across slides.
        """
        slides_map = self._extract_text_exact_logic()
        current_meta = None
        
        for slide_num, text_content in slides_map.items():
            
            # --- CONTEXT DETECTION ---
            # Try to find header in the first few lines of the text block
            lines = text_content.split('\n\n')
            new_meta = None
            
            for line in lines[:5]: # Check top 5 lines
                if "|" in line:
                    # Check against known asset types to confirm it's a header
                    if any(x in line.upper() for x in ["CURRENCY", "ENERGY", "METAL", "GRAIN", "INDEX", "BOND", "SOFT", "MEAT"]):
                        new_meta = self._extract_metadata(line)
                        break
            
            if new_meta:
                current_meta = new_meta

            # --- CONTENT PARSING ---
            if current_meta:
                # Logic Slide Detection
                if "LOGIC BREAKDOWN" in text_content or "1. HIGH VOLATILITY" in text_content:
                    self._parse_logic_content(text_content, current_meta)
                
                # Risk Grid Detection
                if "RISK MANAGEMENT GRID" in text_content or "Entry (Index" in text_content:
                    self._parse_risk_content(text_content, current_meta)

        return pd.DataFrame(self.logic_rows), pd.DataFrame(self.risk_rows)

    def _extract_metadata(self, text_line):
        """Parses 'Product | Asset Class | Timeframe'."""
        clean = re.sub(r"LOGIC BREAKDOWN:|RISK MANAGEMENT GRID:", "", text_line, flags=re.IGNORECASE).strip()
        parts = [p.strip() for p in clean.split('|')]
        if len(parts) >= 3:
            return {
                "Product": parts[0],
                "Asset_Class": parts[1],
                "Timeframe": parts[2],
                "Strategy_ID": f"{parts[0]} ({parts[2]})"
            }
        return None

    def _smart_split_rules(self, text):
        """Splits rules by comma, ignoring brackets."""
        rules = []
        current = []
        depth = 0
        for char in text:
            if char == '(': depth += 1
            elif char == ')': depth = max(0, depth - 1)
            
            if char == ',' and depth == 0:
                rules.append("".join(current).strip())
                current = []
            else:
                current.append(char)
        if current:
            rules.append("".join(current).strip())
        return [r for r in rules if r]

    def _classify_rule(self, indicator):
        if "MA" in indicator.upper(): return "Moving Average"
        return "Filter"

    def _parse_logic_content(self, text, meta):
        """State Machine to parse Logic text."""
        lines = text.split('\n\n') 
        current_regime = None
        current_side = None
        
        for line in lines:
            line = line.strip()
            upper_line = line.upper()
            
            # Regime
            if "1. HIGH VOLATILITY" in upper_line:
                current_regime = "High Vol"
                current_side = None; continue
            elif "2. LOW VOLATILITY" in upper_line:
                current_regime = "Low Vol"
                current_side = None; continue
            elif "3. GENERAL TRAILING" in upper_line:
                current_regime = None; current_side = None; continue

            # Side
            if "[LONG ENTRY]" in upper_line:
                current_side = "Long"; continue
            elif "[SHORT ENTRY]" in upper_line:
                current_side = "Short"; continue

            # Rules
            if current_regime and current_side:
                if "MAs and Filters" in line or "Trailed through" in line or "•" in line:
                    content = line.split(":", 1)[1] if ":" in line else line
                    content = content.replace("•", "").strip()
                    
                    if not content or "No specific rules" in content: continue

                    rules_list = self._smart_split_rules(content)
                    
                    for rule_str in rules_list:
                        match = re.search(r"(.+?)\s*(>=|<=|>|<|=)\s*(.+)", rule_str)
                        indicator = rule_str
                        operator = "Boolean"
                        threshold = None
                        
                        if match:
                            indicator = match.group(1).strip()
                            operator = match.group(2).strip()
                            rhs = match.group(3).strip()
                            try: threshold = float(rhs)
                            except: threshold = None
                        
                        self.logic_rows.append({
                            **meta,
                            "Quadrant": f"{current_side} {current_regime}",
                            "Regime": current_regime,
                            "Side": current_side,
                            "Type": self._classify_rule(indicator),
                            "Indicator": indicator,
                            "Operator": operator,
                            "Threshold": threshold,
                            "Raw_Rule": rule_str
                        })

    def _parse_risk_content(self, text, meta):
        """Parses Risk Grid text."""
        lines = text.split('\n\n')
        current_section = None
        regex = r"Entry \(Index (\d+)\)\s*:\s*Currency Risk:\s*([\d\.]+).*?Target ATR:\s*([\d\.]+).*?Stop ATR:\s*([\d\.]+).*?Risk Reward:\s*([\d\.]+)"
        
        for line in lines:
            line = line.strip()
            upper_line = line.upper()
            
            if "LONG HIGH VOLATILITY" in upper_line: current_section = "Long High Vol"
            elif "LONG LOW VOLATILITY" in upper_line: current_section = "Long Low Vol"
            elif "SHORT HIGH VOLATILITY" in upper_line: current_section = "Short High Vol"
            elif "SHORT LOW VOLATILITY" in upper_line: current_section = "Short Low Vol"
            
            if current_section:
                match = re.search(regex, line)
                if match:
                    self.risk_rows.append({
                        **meta,
                        "Quadrant": current_section,
                        "Entry_Index": int(match.group(1)),
                        "Currency_Risk": float(match.group(2)),
                        "Target_ATR": float(match.group(3)),
                        "Stop_ATR": float(match.group(4)),
                        "RR_Ratio": float(match.group(5))
                    })

# ==========================================
# 2. ANALYSIS LOGIC
# ==========================================

def calculate_detailed_similarity(df):
    """Calculates Jaccard Similarity"""
    pivot = pd.crosstab(df['Strategy_ID'], df['Indicator'])
    if len(pivot) < 2:
        return None, None, None

    dists = pdist(pivot.values, metric='jaccard')
    sim_matrix = 1 - squareform(dists)
    sim_df = pd.DataFrame(sim_matrix, index=pivot.index, columns=pivot.index)

    pairs = []
    for i in range(len(sim_df.columns)):
        for j in range(i + 1, len(sim_df.columns)):
            score = sim_df.iloc[i, j]
            if score > 0.4: 
                strat_a = sim_df.index[i]
                strat_b = sim_df.index[j]
                rules_a = set(df[df['Strategy_ID'] == strat_a]['Indicator'].unique())
                rules_b = set(df[df['Strategy_ID'] == strat_b]['Indicator'].unique())
                
                shared = rules_a.intersection(rules_b)
                unique_a = rules_a - rules_b
                unique_b = rules_b - rules_a
                
                pairs.append({
                    "Strategy A": strat_a,
                    "Strategy B": strat_b,
                    "Score": score,
                    "Shared_Count": len(shared),
                    "Shared_Rules": list(shared),
                    "Unique_A": list(unique_a),
                    "Unique_B": list(unique_b)
                })
    
    pairs_df = pd.DataFrame(pairs)
    if not pairs_df.empty:
        pairs_df = pairs_df.sort_values('Score', ascending=False)
        
    return sim_df, pairs_df, pivot

# ==========================================
# 3. VISUALIZATION MODULES
# ==========================================

def render_logic_dna(df, quadrant_name):
    """
    Module 1: Strategy DNA
    """
    if df.empty:
        st.info(f"No logic data for **{quadrant_name}**.")
        return

    # --- 1. COMPLEXITY HISTOGRAM ---
    rule_counts = df.groupby('Strategy_ID').size().reset_index(name='Rule_Count')
    
    c1, c2 = st.columns([1, 1])
    with c1:
        st.markdown("**1. Complexity Profile**")
        fig_hist = px.histogram(rule_counts, x="Rule_Count", nbins=10, 
                                title="Rules per Strategy", color_discrete_sequence=['#00CC96'])
        st.plotly_chart(fig_hist, use_container_width=True, key=f"comp_hist_{quadrant_name}")

    # --- 2. SIMILARITY ---
    with c2:
        st.markdown("**2. Clone Detector**")
        sim_matrix, pairs_df, pivot = calculate_detailed_similarity(df)
        
        if sim_matrix is not None:
            fig_heat = px.imshow(sim_matrix, x=pivot.index, y=pivot.index,
                                 title="Indicator Overlap Matrix",
                                 color_continuous_scale='RdBu_r', zmin=0, zmax=1)
            st.plotly_chart(fig_heat, use_container_width=True, key=f"sim_heat_{quadrant_name}")
        else:
            st.caption("Not enough data for similarity.")

    # Detailed Pairwise Interpretation
    if sim_matrix is not None and pairs_df is not None and not pairs_df.empty:
        with st.expander("📝 View Detailed Overlap Analysis"):
            for index, row in pairs_df.iterrows():
                score = row['Score']
                s1 = row['Strategy A']
                s2 = row['Strategy B']
                if score > 0.8: icon = "🔴"
                elif score > 0.6: icon = "🟠"
                else: icon = "🟢"
                
                st.markdown(f"**{icon} {s1} vs {s2} ({int(score*100)}% Match)**")
                st.caption(f"Shared Rules: {', '.join(row['Shared_Rules'])}")

    st.markdown("---")
    
    # --- 3. FILTER ANALYSIS ---
    st.markdown("### 3. Indicator Breakdown")
    df_ma = df[df['Type'] == 'Moving Average']
    df_filt = df[df['Type'] == 'Filter']
    
    fc1, fc2 = st.columns(2)
    with fc1:
        st.markdown("**Top Moving Averages**")
        if not df_ma.empty:
            cnt = df_ma['Indicator'].value_counts().reset_index()
            cnt.columns = ['Indicator', 'Count']
            fig = px.bar(cnt.head(8), x='Count', y='Indicator', orientation='h', color_discrete_sequence=['#636EFA'])
            fig.update_layout(yaxis={'categoryorder':'total ascending'})
            st.plotly_chart(fig, use_container_width=True, key=f"ma_{quadrant_name}")
    
    with fc2:
        st.markdown("**Top Value Filters**")
        if not df_filt.empty:
            cnt = df_filt['Indicator'].value_counts().reset_index()
            cnt.columns = ['Indicator', 'Count']
            fig = px.bar(cnt.head(8), x='Count', y='Indicator', orientation='h', color_discrete_sequence=['#EF553B'])
            fig.update_layout(yaxis={'categoryorder':'total ascending'})
            st.plotly_chart(fig, use_container_width=True, key=f"filt_{quadrant_name}")

    # --- 4. THRESHOLD STATS (DETAILED TABLE) ---
    st.markdown("---")
    st.subheader("📊 Detailed Threshold Statistics")
    st.caption("Statistical breakdown of numeric values used in filters for this quadrant.")
    
    num_df = df.dropna(subset=['Threshold'])
    if not num_df.empty:
        # Group by Type, Indicator, Operator
        stats = num_df.groupby(['Type', 'Indicator', 'Operator'])['Threshold'].agg(['count', 'min', 'mean', 'max']).reset_index()
        stats = stats.sort_values(['Type', 'count'], ascending=[True, False])
        
        st.dataframe(
            stats,
            column_config={
                "Type": st.column_config.TextColumn("Category"),
                "Indicator": st.column_config.TextColumn("Indicator"),
                "Operator": st.column_config.TextColumn("Op"),
                "mean": st.column_config.NumberColumn("Avg Value", format="%.4f"),
                "min": st.column_config.NumberColumn("Min", format="%.4f"),
                "max": st.column_config.NumberColumn("Max", format="%.4f"),
                "count": st.column_config.ProgressColumn("Frequency", format="%d", min_value=0, max_value=int(stats['count'].max()))
            },
            use_container_width=True, 
            hide_index=True, 
            key=f"detailed_stats_table_{quadrant_name}"
        )
    else:
        st.info("No numeric thresholds found in this quadrant.")

def render_risk_profile(df, quadrant_name):
    """
    Module 2: Risk Architecture
    """
    if df.empty:
        st.info(f"No risk entries for **{quadrant_name}**.")
        return

    df['BE_WinRate'] = 1 / (1 + df['RR_Ratio'])
    
    r1, r2 = st.columns(2)
    with r1:
        st.markdown("**1. Required Win Rate (Break-Even)**")
        fig_be = px.histogram(df, x="BE_WinRate", nbins=15, 
                              title="Hurdle Rate Distribution",
                              color_discrete_sequence=['#AB63FA'], text_auto=True)
        fig_be.update_layout(xaxis_tickformat=".0%")
        st.plotly_chart(fig_be, use_container_width=True, key=f"be_{quadrant_name}")

    with r2:
        st.markdown("**2. Stop Loss Distribution**")
        fig_stop = px.histogram(df, x="Stop_ATR", nbins=15, 
                                title="Stop Distance (ATR)",
                                color_discrete_sequence=['#FFA15A'])
        st.plotly_chart(fig_stop, use_container_width=True, key=f"stop_{quadrant_name}")

    st.markdown("---")
    st.subheader("📊 Entry Geometry (The 'Jaws')")
    
    # Index Stats
    idx_stats = df.groupby('Entry_Index').agg(
        Avg_Target=('Target_ATR', 'mean'),
        Avg_Stop=('Stop_ATR', 'mean'),
        Avg_RR=('RR_Ratio', 'mean'),
        Count=('Strategy_ID', 'count')
    ).reset_index()
    idx_stats['Entry_Index'] = idx_stats['Entry_Index'].astype(str)

    i1, i2 = st.columns(2)
    with i1:
        st.markdown("**3. Target vs Stop (Visual Edge)**")
        melted = idx_stats.melt(id_vars='Entry_Index', value_vars=['Avg_Target', 'Avg_Stop'], var_name='Metric', value_name='ATR')
        fig_jaws = px.bar(melted, x='Entry_Index', y='ATR', color='Metric', barmode='group',
                          color_discrete_map={'Avg_Target': '#00CC96', 'Avg_Stop': '#EF553B'},
                          title="Gap Analysis")
        st.plotly_chart(fig_jaws, use_container_width=True, key=f"jaws_{quadrant_name}")

    with i2:
        st.markdown("**4. Efficiency (R:R Ratio)**")
        fig_rr = px.bar(idx_stats, x='Entry_Index', y='Avg_RR', text_auto='.2f',
                        title="Avg R:R by Index", color='Avg_RR', color_continuous_scale='Viridis')
        st.plotly_chart(fig_rr, use_container_width=True, key=f"rr_{quadrant_name}")

    st.markdown("**Detailed Index Statistics**")
    st.dataframe(
        idx_stats,
        column_config={
            "Entry_Index": st.column_config.TextColumn("Index"),
            "Count": st.column_config.ProgressColumn("Freq", format="%d", min_value=0, max_value=int(idx_stats['Count'].max())),
            "Avg_Target": st.column_config.NumberColumn("Avg Target", format="%.2f"),
            "Avg_Stop": st.column_config.NumberColumn("Avg Stop", format="%.2f"),
            "Avg_RR": st.column_config.NumberColumn("Avg R:R", format="%.2f"),
        },
        use_container_width=True, hide_index=True, key=f"idx_tbl_{quadrant_name}"
    )

# ==========================================
# 3. MAIN APP CONTROLLER
# ==========================================

def main():
    st.set_page_config(layout="wide", page_title="Hedge Fund Strategy Analytics")
    st.title("📊 Hedge Fund Strategy Analytics")
    st.markdown("---")

    with st.sidebar:
        st.header("1. Data Ingestion")
        uploaded_file = st.file_uploader("Upload 'Strategies_Breakdown.pptx'", type=["pptx"])
        
        df_logic = pd.DataFrame()
        df_risk = pd.DataFrame()
        
        if uploaded_file:
            with st.spinner("Processing..."):
                try:
                    parser = PPTParser(uploaded_file)
                    df_logic, df_risk = parser.parse()
                    if not df_logic.empty:
                        st.success(f"Parsed {df_logic['Strategy_ID'].nunique()} Strategies")
                    else:
                        st.error("Parsing failed. Check file format.")
                except Exception as e:
                    st.error(f"Error: {e}")

        st.divider()
        st.header("2. Analysis Universe")
        
        if df_logic.empty:
            st.info("Upload PPTX to unlock.")
            logic_sub, risk_sub = pd.DataFrame(), pd.DataFrame()
            scope_label = "None"
        else:
            mode = st.radio("Group By:", ["Hierarchy", "Cross-Section", "Custom Selection"])
            logic_sub = df_logic.copy()
            risk_sub = df_risk.copy()
            scope_label = "All"

            # FILTER LOGIC
            if mode == "Hierarchy":
                tfs = sorted(df_logic['Timeframe'].unique())
                sel_tf = st.selectbox("1. Timeframe", tfs)
                assets = sorted(df_logic[df_logic['Timeframe'] == sel_tf]['Asset_Class'].unique())
                sel_asset = st.selectbox("2. Asset Class", assets)
                logic_sub = df_logic[(df_logic['Timeframe'] == sel_tf) & (df_logic['Asset_Class'] == sel_asset)]
                risk_sub = df_risk[(df_risk['Timeframe'] == sel_tf) & (df_risk['Asset_Class'] == sel_asset)]
                scope_label = f"{sel_tf} > {sel_asset}"

            elif mode == "Cross-Section":
                assets = sorted(df_logic['Asset_Class'].unique())
                sel_asset = st.selectbox("Asset Class", assets)
                logic_sub = df_logic[df_logic['Asset_Class'] == sel_asset]
                risk_sub = df_risk[df_risk['Asset_Class'] == sel_asset]
                scope_label = f"{sel_asset} (Global)"

            elif mode == "Custom Selection":
                ids = sorted(df_logic['Strategy_ID'].unique())
                sel_ids = st.multiselect("Strategies", ids, default=ids[:2])
                logic_sub = df_logic[df_logic['Strategy_ID'].isin(sel_ids)]
                risk_sub = df_risk[df_risk['Strategy_ID'].isin(sel_ids)]
                scope_label = "Custom Cluster"

    if logic_sub.empty: return

    with st.container():
        st.subheader(f"Scope: {scope_label}")
        m1, m2, m3, m4 = st.columns(4)
        m1.metric("Active Strategies", logic_sub['Strategy_ID'].nunique())
        m2.metric("Total Rules", len(logic_sub))
        avg_rr = risk_sub['RR_Ratio'].mean() if not risk_sub.empty else 0
        m3.metric("Avg Portfolio R:R", f"{avg_rr:.2f}")
        m4.metric("Avg Required Win Rate", f"{1/(1+avg_rr):.1%}" if avg_rr>0 else "-")
    
    st.divider()

    # TABS: Removed "Comparative Intel"
    tab_dna, tab_risk, tab_raw = st.tabs(["🧬 Strategy DNA", "🛡️ Risk Architecture", "📋 Raw Data"])
    
    quadrants = ["Long High Vol", "Short High Vol", "Long Low Vol", "Short Low Vol"]

    with tab_dna:
        q_tabs = st.tabs(quadrants)
        for i, q in enumerate(quadrants):
            with q_tabs[i]:
                render_logic_dna(logic_sub[logic_sub['Quadrant'] == q], q)

    with tab_risk:
        q_tabs_r = st.tabs(quadrants)
        for i, q in enumerate(quadrants):
            with q_tabs_r[i]:
                render_risk_profile(risk_sub[risk_sub['Quadrant'] == q], q)

    with tab_raw:
        st.subheader("Data Inspector")
        with st.expander("Logic Rules", expanded=True):
            st.dataframe(logic_sub, use_container_width=True)
        with st.expander("Risk Grid"):
            st.dataframe(risk_sub, use_container_width=True)

if __name__ == "__main__":
    main()